from flask import Flask, request, jsonify
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.shapes import MSO_SHAPE_TYPE
import io
import base64
import os
import traceback
from PIL import Image

app = Flask(__name__)

@app.route('/')
def home():
    return "✅ PowerPoint Generator API (cuadros de texto + tablas + logo)."

def reemplazar_marcadores_en_runs(paragraph, nombre_empresa, sector_empresa, detect_logo=False):
    buffer_text = ""
    run_map = []
    logo_detectado = False

    for i, run in enumerate(paragraph.runs):
        buffer_text += run.text
        run_map.append((i, run.text))

    if "{{" not in buffer_text:
        return logo_detectado

    if "{{Logo_Empresa_Cliente}}" in buffer_text:
        logo_detectado = True
        buffer_text = buffer_text.replace("{{Logo_Empresa_Cliente}}", "")

    buffer_text = buffer_text.replace("{{Nombre_Empresa_Cliente}}", nombre_empresa)
    buffer_text = buffer_text.replace("{{Sector_Empresa_Cliente}}", sector_empresa)

    for run in paragraph.runs:
        run.text = ""

    pos = 0
    for i, run in enumerate(paragraph.runs):
        original = run_map[i][1]
        length = len(original)
        if pos >= len(buffer_text):
            break
        run.text = buffer_text[pos:pos+length]
        pos += length

    if pos < len(buffer_text):
        paragraph.add_run(buffer_text[pos:])

    return logo_detectado

@app.route('/generate', methods=['POST'])
def generate_ppt():
    try:
        data = request.get_json()

        nombre_empresa = data.get("Nombre_Empresa_Cliente", "")
        sector_empresa = data.get("Sector_Empresa_Cliente", "")
        logo_data = data.get("Logo_Empresa_Cliente", {}).get("data", "")
        plantilla_data = data.get("Plantilla_Base64", "")

        if not plantilla_data:
            return jsonify({"error": "No se recibió la plantilla (Plantilla_Base64)."}), 400

        plantilla_bytes = base64.b64decode(plantilla_data)
        prs = Presentation(io.BytesIO(plantilla_bytes))

        # Preparar logo
        logo_stream = None
        logo_size = (Inches(1.5), Inches(1.5))
        if logo_data:
            try:
                if isinstance(logo_data, bytes):
                    logo_data = logo_data.decode('utf-8', errors='ignore')
                logo_data = logo_data.replace("\n", "").replace("\r", "")
                image_bytes = base64.b64decode(logo_data)
                image_stream = io.BytesIO(image_bytes)

                img = Image.open(image_stream)
                w, h = img.size
                aspect = w / h
                base_width = Inches(1.5)
                logo_size = (base_width, base_width / aspect)

                image_stream.seek(0)
                logo_stream = image_stream
            except Exception as e:
                return jsonify({"error": f"Logo inválido o corrupto: {str(e)}"}), 400

        # Recorrer todas las diapositivas
        for slide in prs.slides:
            shapes_to_remove = []

            for shape in slide.shapes:
                # Procesar tablas
                if shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                    table = shape.table
                    for row in table.rows:
                        for cell in row.cells:
                            logo_detectado = False
                            for paragraph in cell.text_frame.paragraphs:
                                if reemplazar_marcadores_en_runs(paragraph, nombre_empresa, sector_empresa, detect_logo=True):
                                    logo_detectado = True
                            if logo_detectado and logo_stream:
                                left = shape.left + Inches(0.2)
                                top = shape.top + Inches(0.2)
                                width, height = logo_size
                                logo_stream.seek(0)
                                slide.shapes.add_picture(logo_stream, left, top, width=width, height=height)
                    continue

                # Procesar cuadros de texto
                if not shape.has_text_frame:
                    continue

                text_frame = shape.text_frame
                full_text = "\n".join([p.text for p in text_frame.paragraphs]).strip()

                # Insertar logo si el shape contiene solo el marcador
                if full_text == "{{Logo_Empresa_Cliente}}" and logo_stream:
                    left, top, width, height = shape.left, shape.top, shape.width, shape.height
                    logo_stream.seek(0)
                    slide.shapes.add_picture(logo_stream, left, top, width=width, height=height)
                    shapes_to_remove.append(shape)
                    continue

                # Reemplazar texto en párrafos
                for paragraph in text_frame.paragraphs:
                    reemplazar_marcadores_en_runs(paragraph, nombre_empresa, sector_empresa, detect_logo=True)

            # Eliminar shapes que eran solo el marcador del logo
            for s in shapes_to_remove:
                sp = s._element
                sp.getparent().remove(sp)

        # Guardar resultado
        output = io.BytesIO()
        prs.save(output)
        output.seek(0)
        encoded_output = base64.b64encode(output.read()).decode("utf-8")

        return jsonify({
            "status": "ok",
            "nombre": f"Presentacion_{nombre_empresa}.pptx",
            "file_content": encoded_output
        }), 200

    except Exception as e:
        traceback.print_exc()
        return jsonify({"error": str(e)}), 500

if __name__ == '__main__':
    port = int(os.environ.get("PORT", 10000))
    app.run(host='0.0.0.0', port=port)
